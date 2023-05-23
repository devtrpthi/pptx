import collections.abc
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

# Read content from txt files provided 
with open('sample_slide1_input.txt', 'r') as f:
    slide1_content = f.read()


with open('sample_slide2_input.txt', 'r') as f:
    slide2_content = f.read()


presentation = Presentation()

# Create Slide 1
slide1 = presentation.slides.add_slide(presentation.slide_layouts[1])
slide1.shapes.title.text = "Slide 1"
slide1_content_box = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
#slide1_text_frame = slide1_content_box.text_frame
#slide1_text_frame.text = slide1_content
text_frame = slide1_content_box.text_frame
text_frame.text = slide1_content
text_frame.word_wrap = True
text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
text_frame.text_anchor = MSO_ANCHOR.MIDDLE

# Create Slide 2
slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
slide2.shapes.title.text = "Slide 2"
slide2_content_box = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
slide2_text_frame = slide2_content_box.text_frame
slide2_text_frame.text = slide2_content

# using the given font file to set the fonts
font_file = 'sample_font_file.ttf'
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.file = font_file
                    run.font.name = "Sample Font"
                    #run.font.size = Pt(4)

presentation.save("output.pptx")

